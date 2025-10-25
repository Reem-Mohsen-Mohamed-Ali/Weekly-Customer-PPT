# -*- coding: utf-8 -*-
"""
Created on Sat Oct 25 18:30:04 2025

@author: mreem
"""

# -*- coding: utf-8 -*-
"""
Created on Fri Jun 21 23:40:56 2024

@author: mreem
"""
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from datetime import datetime, timedelta


# Function to extract data from a chart
def extract_chart_data(chart):
    categories = [str(category) for category in chart.plots[0].categories]
    chart_data = []
    for series in chart.series:
        series_data = [value for value in series.values]
        chart_data.append(series_data)
    return categories, chart_data

# Extract data from a specific chart in PowerPoint
def extract_ppt_data(slide_number, chart_title, ppt_file):
    prs = Presentation(ppt_file)
    slide = prs.slides[slide_number]
    for shape in slide.shapes:
        if shape.has_chart and shape.chart.chart_title.text_frame.text.strip() == chart_title:
            return extract_chart_data(shape.chart)
    return None, None

# Function to convert Excel serial date to datetime
def excel_date_to_date(serial):
    try:
        serial = int(serial)
        return (datetime(1899, 12, 30) + timedelta(days=serial)).strftime('%m-%d-%Y')
    except ValueError:
        return serial  # If it's not a valid serial, return as-is

# Parse dates from PowerPoint chart categories
def parse_categories(categories):
    parsed_categories = []
    for category in categories:
        try:
            parsed_categories.append(datetime.strptime(category, '%m-%d-%Y').strftime('%m-%d-%Y'))
        except ValueError:
            parsed_categories.append(excel_date_to_date(category))
    return parsed_categories


def update_ppt_chart(slide_number, chart_title, dates, values, series_name, ppt_file, multi_series=False):

    try:
        prs = Presentation(ppt_file)
        slide = prs.slides[slide_number]

        for shape in slide.shapes:
            if shape.has_chart and shape.chart.chart_title.text_frame.text.strip() == chart_title:
                chart = shape.chart
                chart_data = CategoryChartData()
                chart_data.categories = dates

                if multi_series:
                    # values is a dict: {"Interfreq": [...], "Intrafreq": [...]}
                    for s_name, s_values in values.items():
                        chart_data.add_series(s_name, s_values)
                else:
                    chart_data.add_series(series_name, values)

                chart.replace_data(chart_data)

        prs.save(ppt_file)
        print(f"Updated chart '{chart_title}' on slide {slide_number + 1}")

    except Exception as e:
        print(f"Error updating chart '{chart_title}': {e}")
        
def merge_series(existing_dates, existing_values, new_dates, new_values):
    """
    Merge old and new series by date, removing duplicates.
    New values overwrite old ones for the same date.
    """
    merged = dict(zip(existing_dates, existing_values))  # old data
    merged.update(dict(zip(new_dates, new_values)))     # new data (overwrite)
    merged_dates =  sorted([str(k) for k in merged.keys() if not pd.isna(k)],
    key=lambda x: datetime.strptime(x, "%m-%d-%Y"))
    merged_values = [merged[d] for d in merged_dates]
    return merged_dates, merged_values


###Upper
def main():
    excel_path=r"C:\Users\mreem\Documents\tools\daily\DE_Query 8.xlsx"
    pptx_file = r"C:\Users\mreem\Documents\tools\daily\Delta Overview KPIs.pptx"  # Replace with your PowerPoint file path
    
    prs = Presentation(pptx_file)
    # Update date on the first and twelfth slides
    slides_to_update = [0, 10]  # Slides are 0-indexed


##### Get 5G UPPER Data From Excel #################

    sheet_name = 'Merge1'
    df = pd.read_excel(excel_path, sheet_name=sheet_name)

    # Extract data for LTE CSSR and LTE DCR
    df['Period start time'] = pd.to_datetime(df['Period start time'], errors='coerce')
    dates_NR = df['Period start time'].dt.strftime('%m-%d-%Y').tolist()  # Ensure dates are strings

    
 #### 2G KPI ####################   



    cssr_2G = df['2G Call Setup Success Rate-Speech'].tolist()
    dcr_2G = df['2G Drop Call Rate-Speech_new'].tolist()
    cssr_2G = [float(value) for value in cssr_2G]
    dcr_2G = [float(value) for value in dcr_2G]
    Traffic_2G=df['2G Busy hour Traffic - Speech'].tolist()
    Traffic_2G = [float(value) for value in Traffic_2G] 
    Traffic_Total_2G=df['2G Data'].tolist()
    Traffic_Total_2G = [float(value) for value in Traffic_Total_2G] 
    
 #### 3G KPI ####################   


    cssr_3G_V = df['3G CAll SETUP SUCCESS RATE - SPEECH'].tolist()
    cssr_3G_D = df['3G CALL SETUP SUCCESS RATE - HSDPA'].tolist()
    dcr_3G_V = df['3G VOICE_DCR'].tolist()
    dcr_3G_D = df['DCR_HSDPA'].tolist()

    cssr_3G_V = [float(value) for value in cssr_3G_V]
    dcr_3G_V = [float(value) for value in dcr_3G_V]
    cssr_3G_D = [float(value) for value in cssr_3G_D]
    dcr_3G_D = [float(value) for value in dcr_3G_D]
    Traffic_3G=df['TRAFFIC_SPEECH2'].tolist()
    Traffic_3G = [float(value) for value in Traffic_3G]
    Total_3G_Traffic = df['3G Data'].tolist()
    Total_3G_Traffic = [float(value) for value in Total_3G_Traffic]

    
    

 #### 4G KPI ####################   

    cssr_4G_D = df['4G/LTE CALL SETUP SUCCESS RATE'].tolist()
    cssr_4G_V = df['CSSR_VOLTE'].tolist()
    dcr_4G_D = df['4G/LTE DROP CALL RATE (ALL)'].tolist()
    dcr_4G_V = df['DCR_VOLTE'].tolist()  
    
    Traffic_4G_Data = df['4G Data'].tolist()   
    Traffic_4G_VO = df['VoLTE total traffic'].tolist() 
    Traffic_ORG= df['Orange Traffic'].tolist() 
    Traffic_WE= df['WE Traffic'].tolist() 
    Traffic_3G_PEN= df['3G Penetration'].tolist() 
    Traffic_VOLTE_PEN= df['VoLTE Penetration'].tolist() 
    Traffic_4G_PEN= df['4G Penetration'].tolist() 
    Traffic_WE_PEN= df['QCI9 Share'].tolist() 
    
    cssr_4G_D = [float(value) for value in cssr_4G_D]  
    cssr_4G_V = [float(value) for value in cssr_4G_V]  
    dcr_4G_D = [float(value) for value in dcr_4G_D]  
    dcr_4G_V = [float(value) for value in dcr_4G_V]     
    Traffic_4G_Data = [float(value) for value in Traffic_4G_Data]
    Traffic_4G_VO = [float(value) for value in Traffic_4G_VO]        
    Traffic_ORG = [float(value) for value in Traffic_ORG]    
    Traffic_WE = [float(value) for value in Traffic_WE]  
    Traffic_3G_PEN = [float(value) for value in Traffic_3G_PEN]      
    Traffic_VOLTE_PEN = [float(value) for value in Traffic_VOLTE_PEN]       
    Traffic_4G_PEN = [float(value) for value in Traffic_4G_PEN]       
    Traffic_WE_PEN = [float(value) for value in Traffic_WE_PEN]     

    

    Traffic_5G_DL_CO = df['5G leg'].tolist() 
    Traffic_5G_DL_CO = [float(value) for value in Traffic_5G_DL_CO]
    Traffic_x2 = df['X2 total'].tolist() 
    Traffic_x2 = [float(value) for value in Traffic_x2]
    Traffic_pdcp = df['PDCP 4G total'].tolist() 
    Traffic_pdcp = [float(value) for value in Traffic_pdcp]  
    Traffic_5GCO = df['5G Collocated'].tolist() 
    Traffic_5GCO = [float(value) for value in Traffic_5GCO]  
    
    ############################ 5G KPIs ###################
    ##### Get 5G UPPER Data From Excel #################
    shhet_5g=r"C:\Users\mreem\Documents\tools\daily\WeeklySlides_Govs-nsn2gnpo-2025_10_24-18_03_56__215 1.xlsx"
    
    sheet_name = '5G001_Govs'
    df = pd.read_excel(shhet_5g, sheet_name=sheet_name)
    filtered_df = df[df['WS_NAME'] == '5G_Colloc_5G_pilot_delta']
    filtered_df['Period start time'] = pd.to_datetime(filtered_df['Period start time'], errors='coerce')
    # Extract data for LTE CSSR and LTE DCR
    dates_NR_2 = filtered_df['Period start time'].dt.strftime('%Y-%m-%d').tolist()  # Ensure dates are strings
    NR_cssr_data = filtered_df['NSA call access'].tolist()
    
    NR_HOSR_INTER_data = filtered_df['IntergNB HO SR NSA'].tolist()
    NR_HOSR_INTRA_data = filtered_df['Inafreq inaDU PSC chg tot SR'].tolist()
    MAX_USER_THR_DL_data = filtered_df['Max MAC SDU Cell Thr DL DTCH'].tolist()
    MAX_USER_THR_UL_data = filtered_df['Max MAC SDU Cell Thr UL DTCH'].tolist()
    USER_THR_UL_data = filtered_df['Avg MAC user thp UL'].tolist()
    NR_USER_AVG = filtered_df['NSA Avg nr user'].tolist()
    USER_THR_DL_data = filtered_df['Avg MAC user thp DL exc DRX sleep'].tolist()
    USER_THR_DL_data = [float(value) for value in USER_THR_DL_data]

    dates_NR_2= [datetime.strptime(date, '%Y-%m-%d').strftime('%m-%d-%Y') for date in dates_NR_2]
    # Convert float values in the lists if they are not already floats
    NR_cssr_data = [float(value) for value in NR_cssr_data]

    NR_HOSR_INTER_data = [float(value) for value in NR_HOSR_INTER_data]
    NR_HOSR_INTRA_data = [float(value) for value in NR_HOSR_INTRA_data]
    MAX_USER_THR_DL_data = [float(value) for value in MAX_USER_THR_DL_data]
    MAX_USER_THR_UL_data = [float(value) for value in MAX_USER_THR_UL_data]
    USER_THR_UL_data = [float(value) for value in USER_THR_UL_data]
    NR_USER_AVG = [float(value) for value in NR_USER_AVG]

#### 5G Rest of KPI ####################
    sheet_name_5g = '5G_DCR'
    df_5g = pd.read_excel(shhet_5g, sheet_name=sheet_name_5g)
    filtered_df_5g = df_5g[df_5g['WS_NAME'] == '5G_Colloc_5G_pilot_delta']
    NR_dcr_data = filtered_df_5g['NSA SgNB t abn rel R excl X2 rst'].tolist()
    NR_dcr_data = [float(value) for value in NR_dcr_data]

########## 1st slide ########################################################################################################################################
    slide_number_Voice_Traffic = 1
    chart_title_Voice_Traffic = "Voice Traffic Evolution"
    chart_title_Data_Traffic = "Data Traffic Evolution"
    chart_title_WE_Traffic = "National Roaming Traffic Share"
    chart_title_NR_Traffic = "5G Collocated Traffic Share"

########## Update Voice_Traffic KPIS (same style as CSSR/DCR)###################
    categories_Voice_Traffic, existing_Voice_Traffic = extract_ppt_data(slide_number_Voice_Traffic, chart_title_Voice_Traffic, pptx_file)
    parsed_categories_Voice_Traffic = parse_categories(categories_Voice_Traffic)
    existing_Voice_Traffic[0] = [float(v) for v in existing_Voice_Traffic[0]] 
    existing_Voice_Traffic[1] = [float(v) for v in existing_Voice_Traffic[1]] 
    existing_Voice_Traffic[2] = [float(v) for v in existing_Voice_Traffic[2]] 
    existing_Voice_Traffic[3] = [float(v) for v in existing_Voice_Traffic[3]]
    existing_Voice_Traffic[4] = [float(v) for v in existing_Voice_Traffic[4]]
    combined_values_Voice_Traffic = {
    "2G Busy hour Traffic - Speech": existing_Voice_Traffic[0] + Traffic_2G,
    "Traffic Voice": existing_Voice_Traffic[1] + Traffic_3G,
   "VoLTE total traffic": existing_Voice_Traffic[2] + Traffic_4G_VO,
   "VoLTE Penetration": existing_Voice_Traffic[3] + Traffic_VOLTE_PEN,
   "3G Penetration": existing_Voice_Traffic[3] + Traffic_3G_PEN
   } 
        
    merged_dates_Voice_Traffic_2G, merged_values_Voice_Traffic_2G = merge_series(parsed_categories_Voice_Traffic, existing_Voice_Traffic[0], dates_NR, Traffic_2G)
    merged_dates_Voice_Traffic_3G, merged_values_Voice_Traffic_3G = merge_series(parsed_categories_Voice_Traffic, existing_Voice_Traffic[1], dates_NR, Traffic_3G)
    merged_dates_Voice_Traffic_4G, merged_values_Voice_Traffic_4G = merge_series(parsed_categories_Voice_Traffic, existing_Voice_Traffic[2], dates_NR, Traffic_4G_VO)
    merged_dates_Voice_Traffic_VP, merged_values_Voice_Traffic_VP = merge_series(parsed_categories_Voice_Traffic, existing_Voice_Traffic[3], dates_NR, Traffic_VOLTE_PEN)
    merged_dates_Voice_Traffic_3GP, merged_values_Voice_Traffic_3GP = merge_series(parsed_categories_Voice_Traffic, existing_Voice_Traffic[4], dates_NR, Traffic_3G_PEN)    
    combined_values_Voice_Traffic = {
    "2G Busy hour Traffic - Speech": merged_values_Voice_Traffic_2G,
    "Traffic Voice": merged_values_Voice_Traffic_3G,
   "VoLTE total traffic": merged_values_Voice_Traffic_4G,
   "VoLTE Penetration": merged_values_Voice_Traffic_VP,
   "3G Penetration": merged_values_Voice_Traffic_3GP} 
        
        
    update_ppt_chart(
        slide_number_Voice_Traffic,
        chart_title_Voice_Traffic,
        merged_dates_Voice_Traffic_2G,
        combined_values_Voice_Traffic,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True) 

############ Update Data_Traffic KPIS (same style as CSSR/DCR)##################
    categories_Data_Traffic, existing_Data_Traffic = extract_ppt_data(slide_number_Voice_Traffic, chart_title_Data_Traffic, pptx_file)
    parsed_categories_Data_Traffic = parse_categories(categories_Data_Traffic)
    existing_Data_Traffic[0] = [float(v) for v in existing_Data_Traffic[0]] 
    existing_Data_Traffic[1] = [float(v) for v in existing_Data_Traffic[1]] 
    existing_Data_Traffic[2] = [float(v) for v in existing_Data_Traffic[2]] 
    existing_Data_Traffic[3] = [float(v) for v in existing_Data_Traffic[3]] 
    combined_values_Data_Traffic = {
    "Sum of 2G Data": existing_Data_Traffic[0] + Traffic_Total_2G,
    "Sum of 3G Data": existing_Data_Traffic[1] + Total_3G_Traffic,
   "Sum of 4G Data": existing_Data_Traffic[2] + Traffic_4G_Data,
   "LTE Penetration": existing_Data_Traffic[3] + Traffic_4G_PEN} 
        
    merged_dates_Data_Traffic_2G_d, merged_values_Data_Traffic_2G_d = merge_series(parsed_categories_Data_Traffic, existing_Data_Traffic[0], dates_NR, Traffic_Total_2G)
    merged_dates_Data_Traffic_3G_d, merged_values_Data_Traffic_3G_d = merge_series(parsed_categories_Data_Traffic, existing_Data_Traffic[1], dates_NR, Total_3G_Traffic)
    merged_dates_Data_Traffic_4G_d, merged_values_Data_Traffic_4G_d = merge_series(parsed_categories_Data_Traffic, existing_Data_Traffic[2], dates_NR, Traffic_4G_Data)
    merged_dates_Data_Traffic_tP_d, merged_values_Data_Traffic_tP_d = merge_series(parsed_categories_Data_Traffic, existing_Data_Traffic[3], dates_NR, Traffic_4G_PEN)
  
    combined_values_Data_Traffic = {
    "Sum of 2G Data": merged_values_Data_Traffic_2G_d,
    "Sum of 3G Data": merged_values_Data_Traffic_3G_d,
   "Sum of 4G Data": merged_values_Data_Traffic_4G_d,
   "LTE Penetration": merged_values_Data_Traffic_tP_d} 

        
        
    update_ppt_chart(
        slide_number_Voice_Traffic,
        chart_title_Data_Traffic,
        merged_dates_Data_Traffic_2G_d,
        combined_values_Data_Traffic,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True) 
    
    
############ Update WE_Traffic KPIS (same style as CSSR/DCR)###################
    categories_WE_Traffic, existing_WE_Traffic = extract_ppt_data(slide_number_Voice_Traffic, chart_title_WE_Traffic, pptx_file)
    parsed_categories_WE_Traffic = parse_categories(categories_WE_Traffic)
    existing_WE_Traffic[0] = [float(v) for v in existing_WE_Traffic[0]] 
    existing_WE_Traffic[1] = [float(v) for v in existing_WE_Traffic[1]] 
    existing_WE_Traffic[2] = [float(v) for v in existing_WE_Traffic[2]] 
    combined_values_WE_Traffic = {
    "Orange Traffic (QCI8)": existing_WE_Traffic[0] + Traffic_ORG,
    "WE Traffic (QCI9)": existing_WE_Traffic[1] + Traffic_WE,
   "WE traffic Share": existing_WE_Traffic[2] + Traffic_WE_PEN} 
        
    merged_dates_WE_Traffic_ORG, merged_values_WE_Traffic_ORG = merge_series(parsed_categories_WE_Traffic, existing_WE_Traffic[0], dates_NR, Traffic_ORG)
    merged_dates_WE_Traffic_WE, merged_values_WE_Traffic_WE = merge_series(parsed_categories_WE_Traffic, existing_WE_Traffic[1], dates_NR, Traffic_WE)
    merged_dates_WE_Traffic_WP, merged_values_WE_Traffic_WP = merge_series(parsed_categories_WE_Traffic, existing_WE_Traffic[2], dates_NR, Traffic_WE_PEN)

  
    combined_values_WE_Traffic = {
    "Orange Traffic (QCI8)": merged_values_WE_Traffic_ORG,
    "WE Traffic (QCI9)": merged_values_WE_Traffic_WE,
   "WE traffic Share": merged_values_WE_Traffic_WP} 

        
        
    update_ppt_chart(
        slide_number_Voice_Traffic,
        chart_title_WE_Traffic,
        merged_dates_WE_Traffic_ORG,
        combined_values_WE_Traffic,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True) 
  
############# Update NR_Traffic KPIS (same style as CSSR/DCR)###################
    categories_NR_Traffic, existing_NR_Traffic = extract_ppt_data(slide_number_Voice_Traffic, chart_title_NR_Traffic, pptx_file)
    parsed_categories_NR_Traffic = parse_categories(categories_NR_Traffic)
    existing_NR_Traffic[0] = [float(v) for v in existing_NR_Traffic[0]] 
    existing_NR_Traffic[1] = [float(v) for v in existing_NR_Traffic[1]] 
    existing_NR_Traffic[2] = [float(v) for v in existing_NR_Traffic[2]]
    existing_NR_Traffic[3] = [float(v) for v in existing_NR_Traffic[3]]

    combined_values_NR_Traffic = {
    "PDCP 4G Total": existing_NR_Traffic[0] + Traffic_pdcp,
    "5G leg": existing_NR_Traffic[1] + Traffic_5G_DL_CO,
   "X2 total": existing_NR_Traffic[2] + Traffic_x2,
   "5G Collocated": existing_NR_Traffic[3] + Traffic_5GCO } 
        
    merged_dates_NR_Traffic_LTE, merged_values_NR_Traffic_LTE = merge_series(parsed_categories_NR_Traffic, existing_NR_Traffic[0], dates_NR, Traffic_pdcp)
    merged_dates_NR_Traffic_NR, merged_values_NR_Traffic_NR = merge_series(parsed_categories_NR_Traffic, existing_NR_Traffic[1], dates_NR, Traffic_5G_DL_CO)
    merged_dates_NR_Traffic_NR, merged_values_NR_Traffic_x2 = merge_series(parsed_categories_NR_Traffic, existing_NR_Traffic[2], dates_NR, Traffic_x2)
    merged_dates_NR_Traffic_NP, merged_values_NR_Traffic_NP = merge_series(parsed_categories_NR_Traffic, existing_NR_Traffic[3], dates_NR, Traffic_5GCO)

  
    combined_values_NR_Traffic = {
    "PDCP 4G Total": merged_values_NR_Traffic_LTE,
    "5G leg": merged_values_NR_Traffic_NR,
   "X2 total": merged_values_NR_Traffic_x2,
   "5G Collocated":merged_values_NR_Traffic_NP} 

        
        
    update_ppt_chart(
        slide_number_Voice_Traffic,
        chart_title_NR_Traffic,
        merged_dates_NR_Traffic_LTE,
        combined_values_NR_Traffic,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)   
  
########## 4th slide ##########################################################################################################################
    slide_number_LTE_Traffic = 2
    chart_title_LTE_CSSR = "4G CSSR"
    chart_title_LTE_DCR = "4G DCR"
    chart_title_LTE_Traffic = "LTE & VoLTE Traffic"

############# Update LTE_CSSR KPIS (same style as CSSR/DCR)###################
    categories_LTE_CSSR, existing_LTE_CSSR = extract_ppt_data(slide_number_LTE_Traffic, chart_title_LTE_CSSR, pptx_file)
    parsed_categories_LTE_CSSR = parse_categories(categories_LTE_CSSR)
    existing_LTE_CSSR[0] = [float(v) for v in existing_LTE_CSSR[0]] 
    existing_LTE_CSSR[1] = [float(v) for v in existing_LTE_CSSR[1]]  
    
    combined_values_LTE_CSSR = {
    "LTE_CSSR_Orange V5.2": existing_LTE_CSSR[0] + cssr_4G_D,
    "Orange CSSR VOLTE 2": existing_LTE_CSSR[1] + cssr_4G_V} 
        
    merged_dates_LTE_CSSR_D, merged_values_LTE_CSSR_D = merge_series(parsed_categories_LTE_CSSR, existing_LTE_CSSR[0], dates_NR, cssr_4G_D)
    merged_dates_LTE_CSSRc_3G, merged_values_LTE_CSSR_V = merge_series(parsed_categories_LTE_CSSR, existing_LTE_CSSR[1], dates_NR, cssr_4G_V)
    
    combined_values_LTE_CSSR = {
    "LTE_CSSR_Orange V5.2": merged_values_LTE_CSSR_D,
    "Orange CSSR VOLTE 2": merged_values_LTE_CSSR_V} 
        
        
    update_ppt_chart(
        slide_number_LTE_Traffic,
        chart_title_LTE_CSSR,
        merged_dates_LTE_CSSR_D,
        combined_values_LTE_CSSR,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)    
    
############# Update LTE_DCR KPIS (same style as CSSR/DCR)###################
    categories_LTE_DCR, existing_LTE_DCR = extract_ppt_data(slide_number_LTE_Traffic, chart_title_LTE_DCR, pptx_file)
    parsed_categories_LTE_DCR = parse_categories(categories_LTE_DCR)
    existing_LTE_DCR[0] = [float(v) for v in existing_LTE_DCR[0]] 
    existing_LTE_DCR[1] = [float(v) for v in existing_LTE_DCR[1]]  
    
    combined_values_LTE_DCR = {
    "LTE_DCR_v5_RRC_TA_filter": existing_LTE_DCR[0] + dcr_4G_D,
    "Orange VOLTE DCR V5.2": existing_LTE_DCR[1] + dcr_4G_V} 
        
    merged_dates_LTE_DCR_D, merged_values_LTE_DCR_D = merge_series(parsed_categories_LTE_DCR, existing_LTE_DCR[0], dates_NR, dcr_4G_D)
    merged_dates_LTE_DCR_V, merged_values_LTE_DCR_V = merge_series(parsed_categories_LTE_DCR, existing_LTE_DCR[1], dates_NR, dcr_4G_V)
    
    combined_values_LTE_DCR = {
    "LTE_DCR_v5_RRC_TA_filter": merged_values_LTE_DCR_D,
    "Orange VOLTE DCR V5.2": merged_values_LTE_DCR_V} 
        
        
    update_ppt_chart(
        slide_number_LTE_Traffic,
        chart_title_LTE_DCR,
        merged_dates_LTE_DCR_D,
        combined_values_LTE_DCR,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)     
    
############# Update LTE_traff KPIS (same style as CSSR/DCR)###################
    categories_LTE_traff, existing_LTE_traff = extract_ppt_data(slide_number_LTE_Traffic, chart_title_LTE_Traffic, pptx_file)
    parsed_categories_LTE_traff = parse_categories(categories_LTE_traff)
    existing_LTE_traff[0] = [float(v) for v in existing_LTE_traff[0]] 
    existing_LTE_traff[1] = [float(v) for v in existing_LTE_traff[1]]  
    
    combined_values_LTE_traff = {
    "4G Data Traffic": existing_LTE_traff[0] + Traffic_4G_Data,
    "VoLTE total traffic": existing_LTE_traff[1] + Traffic_4G_VO} 
        
    merged_dates_LTE_traff_D, merged_values_LTE_traff_D = merge_series(parsed_categories_LTE_traff, existing_LTE_traff[0], dates_NR, Traffic_4G_Data)
    merged_dates_LTE_traff_V, merged_values_LTE_traff_V = merge_series(parsed_categories_LTE_traff, existing_LTE_traff[1], dates_NR, Traffic_4G_VO)
    
    combined_values_LTE_traff = {
    "4G Data Traffic": merged_values_LTE_traff_D,
    "VoLTE total traffic": merged_values_LTE_traff_V} 
        
        
    update_ppt_chart(
        slide_number_LTE_Traffic,
        chart_title_LTE_Traffic,
        merged_dates_LTE_traff_D,
        combined_values_LTE_traff,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)    
    
########## 6th slide ##########################################################################################################################
    slide_number_3G_Traffic = 3
    chart_title_3G_CSSR = "3G CSSR"
    chart_title_3G_DCR = "3G DCR"
    chart_title_2G_CSSR = "2G CSSR"
    chart_title_2G_DCR = "2G DCR"


############# Update 3G_CSSR KPIS (same style as CSSR/DCR)###################
    categories_3G_CSSR, existing_3G_CSSR = extract_ppt_data(slide_number_3G_Traffic, chart_title_3G_CSSR, pptx_file)
    parsed_categories_3G_CSSR = parse_categories(categories_3G_CSSR)
    existing_3G_CSSR[0] = [float(v) for v in existing_3G_CSSR[0]] 
    existing_3G_CSSR[1] = [float(v) for v in existing_3G_CSSR[1]]  
    
    combined_values_3G_CSSR = {
    "3G CAll SETUP SUCCESS RATE - SPEECH": existing_3G_CSSR[0] + cssr_3G_V,
    "3G CALL SETUP SUCCESS RATE - HSDPA": existing_3G_CSSR[1] + cssr_3G_D} 
        
    merged_dates_3G_CSSR_V, merged_values_3G_CSSR_V = merge_series(parsed_categories_3G_CSSR, existing_3G_CSSR[0], dates_NR, cssr_3G_V)
    merged_dates_3G_CSSR_D, merged_values_3G_CSSR_D = merge_series(parsed_categories_3G_CSSR, existing_3G_CSSR[1], dates_NR, cssr_3G_D)
    
    combined_values_3G_CSSR = {
    "3G CAll SETUP SUCCESS RATE - SPEECH": merged_values_3G_CSSR_V,
    "3G CALL SETUP SUCCESS RATE - HSDPA": merged_values_3G_CSSR_D} 
        
        
    update_ppt_chart(
        slide_number_3G_Traffic,
        chart_title_3G_CSSR,
        merged_dates_3G_CSSR_D,
        combined_values_3G_CSSR,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)    
    

############# Update 3G_DCR KPIS (same style as CSSR/DCR)###################
    categories_3G_DCR, existing_3G_DCR = extract_ppt_data(slide_number_3G_Traffic, chart_title_3G_DCR, pptx_file)
    parsed_categories_3G_DCR = parse_categories(categories_3G_DCR)
    existing_3G_DCR[0] = [float(v) for v in existing_3G_DCR[0]] 
    existing_3G_DCR[1] = [float(v) for v in existing_3G_DCR[1]]  
    
    combined_values_3G_DCR = {
    "voicedcr3g": existing_3G_DCR[0] + dcr_3G_V,
    "HSDPA_DCR_Weekly": existing_3G_DCR[1] + dcr_3G_D} 
        
    merged_dates_3G_DCR_V, merged_values_3G_DCR_V = merge_series(parsed_categories_3G_DCR, existing_3G_DCR[0], dates_NR, dcr_3G_V)
    merged_dates_3G_DCR_D, merged_values_3G_DCR_D = merge_series(parsed_categories_3G_DCR, existing_3G_DCR[1], dates_NR, dcr_3G_D)
    
    combined_values_3G_DCR = {
    "voicedcr3g": merged_values_3G_DCR_V,
    "HSDPA_DCR_Weekly": merged_values_3G_DCR_D} 
        
        
    update_ppt_chart(
        slide_number_3G_Traffic,
        chart_title_3G_DCR,
        merged_dates_3G_DCR_D,
        combined_values_3G_DCR,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)     
    
############# Update 2G_CSSR KPIS (same style as CSSR/DCR)###################
    categories_2G_CSSR, existing_2G_CSSR = extract_ppt_data(slide_number_3G_Traffic, chart_title_2G_CSSR, pptx_file)
    parsed_categories_2G_CSSR = parse_categories(categories_2G_CSSR)
    existing_2G_CSSR[0] = [float(v) for v in existing_2G_CSSR[0]] 

    
    combined_values_2G_CSSR = {
    "GSMCSSRVoicenew": existing_2G_CSSR[0] + cssr_2G}
        
    merged_dates_2G_CSSR_V, merged_values_2G_CSSR_V = merge_series(parsed_categories_2G_CSSR, existing_2G_CSSR[0], dates_NR, cssr_2G)

    
    combined_values_2G_CSSR = {
    "GSMCSSRVoicenew": merged_values_2G_CSSR_V}
        
        
    update_ppt_chart(
        slide_number_3G_Traffic,
        chart_title_2G_CSSR,
        merged_dates_2G_CSSR_V,
        combined_values_2G_CSSR,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)    
    

############# Update 2G_DCR KPIS (same style as CSSR/DCR)###################
    categories_2G_DCR, existing_2G_DCR = extract_ppt_data(slide_number_3G_Traffic, chart_title_2G_DCR, pptx_file)
    parsed_categories_2G_DCR = parse_categories(categories_2G_DCR)
    existing_2G_DCR[0] = [float(v) for v in existing_2G_DCR[0]] 

    
    combined_values_2G_DCR = {
    "GSMCSSRVoicenew": existing_2G_DCR[0] + dcr_2G}
        
    merged_dates_2G_DCR_V, merged_values_2G_DCR_V = merge_series(parsed_categories_2G_DCR, existing_2G_DCR[0], dates_NR, dcr_2G)

    
    combined_values_2G_DCR = {
    "GSMCSSRVoicenew": merged_values_2G_DCR_V}
        
        
    update_ppt_chart(
        slide_number_3G_Traffic,
        chart_title_2G_DCR,
        merged_dates_2G_DCR_V,
        combined_values_2G_DCR,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)
########## 6th slide ##########################################################################################################################
    slide_number_5G_Traffic = 4
    chart_title_5G_CSSR = "CSSR"
    chart_title_5G_DCR = "DCR"
    chart_title_5G_HOSR = "HOSR"
    chart_title_NSA_USER = "NSA Users"
    chart_title_Avg_Tput = "Avg Tput"
    chart_title_Max_Tput = "Max Tput"
    
############# Update 5G_CSSR KPIS (same style as CSSR/DCR)###################
    categories_5G_CSSR, existing_5G_CSSR = extract_ppt_data(slide_number_5G_Traffic, chart_title_5G_CSSR, pptx_file)
    parsed_categories_5G_CSSR = parse_categories(categories_5G_CSSR)
    existing_5G_CSSR[0] = [float(v) for v in existing_5G_CSSR[0]] 
    print(categories_5G_CSSR)
    
    combined_values_5G_CSSR = {
    "NSA call access": existing_5G_CSSR[0] + NR_cssr_data}
        
    merged_dates_5G_CSSR_V, merged_values_5G_CSSR_V = merge_series(parsed_categories_5G_CSSR, existing_5G_CSSR[0], dates_NR_2, NR_cssr_data)

    
    combined_values_5G_CSSR = {
    "NSA call access": merged_values_5G_CSSR_V}
        
        
    update_ppt_chart(
        slide_number_5G_Traffic,
        chart_title_5G_CSSR,
        merged_dates_5G_CSSR_V,
        combined_values_5G_CSSR,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)    
  

############# Update 5G_DCR KPIS (same style as CSSR/DCR)###################
    categories_5G_DCR, existing_5G_DCR = extract_ppt_data(slide_number_5G_Traffic, chart_title_5G_DCR, pptx_file)
    parsed_categories_5G_DCR = parse_categories(categories_5G_DCR)
    existing_5G_DCR[0] = [float(v) for v in existing_5G_DCR[0]] 

    
    combined_values_5G_DCR = {
    "5G NSA DROP CALL RATE NRBTS": existing_5G_DCR[0] + NR_dcr_data}
        
    merged_dates_5G_DCR_V, merged_values_5G_DCR_V = merge_series(parsed_categories_5G_DCR, existing_5G_DCR[0], dates_NR_2, NR_dcr_data)

    
    combined_values_5G_DCR = {
    "5G NSA DROP CALL RATE NRBTS": merged_values_5G_DCR_V}
        
        
    update_ppt_chart(
        slide_number_5G_Traffic,
        chart_title_5G_DCR,
        merged_dates_5G_DCR_V,
        combined_values_5G_DCR,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)    
    
############# Update 5G_DCR KPIS (same style as CSSR/DCR)###################
    categories_NSA_USER, existing_NSA_USER = extract_ppt_data(slide_number_5G_Traffic, chart_title_NSA_USER, pptx_file)
    parsed_categories_NSA_USER = parse_categories(categories_NSA_USER)
    existing_NSA_USER[0] = [float(v) for v in existing_NSA_USER[0]] 

    
    combined_values_NSA_USER = {
    "NSA Avg nr user": existing_NSA_USER[0] + NR_USER_AVG}
        
    merged_dates_NSA_USER_V, merged_values_NSA_USER_V = merge_series(parsed_categories_NSA_USER, existing_NSA_USER[0], dates_NR_2, NR_USER_AVG)

    
    combined_values_NSA_USER = {
    "NSA Avg nr user": merged_values_NSA_USER_V}
        
 
    update_ppt_chart(
        slide_number_5G_Traffic,
        chart_title_NSA_USER,
        merged_dates_NSA_USER_V,
        combined_values_NSA_USER,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)    
    
############# Update HOSR KPIS (same style as CSSR/DCR)###################
    categories_NR_HOSR, existing_NR_HOSR = extract_ppt_data(slide_number_5G_Traffic, chart_title_5G_HOSR, pptx_file)
    parsed_categories_NR_HOSR = parse_categories(categories_NR_HOSR)
    existing_NR_HOSR[0] = [float(v) for v in existing_NR_HOSR[0]] 
    existing_NR_HOSR[1] = [float(v) for v in existing_NR_HOSR[1]]  
    
    combined_values_NR_HOSR = {
    "IntergNB HO SR NSA": existing_NR_HOSR[0] + NR_HOSR_INTER_data,
    "Inafreq inaDU PSC chg tot SR": existing_NR_HOSR[1] + NR_HOSR_INTRA_data} 
        
    merged_dates_NR_HOSR_inter, merged_values_NR_HOSR_inter = merge_series(parsed_categories_NR_HOSR, existing_NR_HOSR[0], dates_NR_2, NR_HOSR_INTER_data)
    merged_dates_NR_HOSR_intra, merged_values_NR_HOSR_intra = merge_series(parsed_categories_NR_HOSR, existing_NR_HOSR[1], dates_NR_2, NR_HOSR_INTRA_data)
    
    combined_values_NR_HOSR = {
    "IntergNB HO SR NSA": merged_values_NR_HOSR_inter,
    "Inafreq inaDU PSC chg tot SR": merged_values_NR_HOSR_intra} 
        
        
    update_ppt_chart(
        slide_number_5G_Traffic,
        chart_title_5G_HOSR,
        merged_dates_NR_HOSR_inter,
        combined_values_NR_HOSR,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)     

############# Update Max Tput KPIS (same style as CSSR/DCR)###################

    categories_max_tput, existing_max_tput = extract_ppt_data(slide_number_5G_Traffic, chart_title_Max_Tput, pptx_file)
    parsed_categories_max_tput = parse_categories(categories_max_tput)
    existing_max_tput[0] = [float(v) for v in existing_max_tput[0]] 
    existing_max_tput[1] = [float(v) for v in existing_max_tput[1]]  
    
    combined_values_max_tput = {
    "Max MAC SDU Cell Thr DL DTCH": existing_max_tput[0] + MAX_USER_THR_DL_data,
    "Max MAC SDU Cell Thr UL DTCH": existing_max_tput[1] + MAX_USER_THR_UL_data} 
        
    merged_dates_max_tput_DL, merged_values_max_tput_DL = merge_series(parsed_categories_max_tput, existing_max_tput[0], dates_NR_2, MAX_USER_THR_DL_data)
    merged_dates_max_tput_UL, merged_values_max_tput_UL = merge_series(parsed_categories_max_tput, existing_max_tput[1], dates_NR_2, MAX_USER_THR_UL_data)
    
    combined_values_max_tput = {
    "Max MAC SDU Cell Thr DL DTCH": merged_values_max_tput_DL,
    "Max MAC SDU Cell Thr UL DTCH": merged_values_max_tput_UL} 
        
        
    update_ppt_chart(
        slide_number_5G_Traffic,
        chart_title_Max_Tput,
        merged_dates_max_tput_DL,
        combined_values_max_tput,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)   
    
############# Update Avg Tput KPIS (same style as CSSR/DCR)###################

    categories_Avg_Tput, existing_Avg_Tput = extract_ppt_data(slide_number_5G_Traffic, chart_title_Avg_Tput, pptx_file)
    parsed_categories_Avg_Tput = parse_categories(categories_Avg_Tput)
    existing_Avg_Tput[0] = [float(v) for v in existing_Avg_Tput[0]] 
    existing_Avg_Tput[1] = [float(v) for v in existing_Avg_Tput[1]]  
    
    combined_values_Avg_Tput = {
    "Avg MAC user thp DL exc DRX sleep": existing_Avg_Tput[0] + USER_THR_DL_data,
    "Avg MAC user thp UL": existing_Avg_Tput[1] + USER_THR_UL_data} 
        
    merged_dates_Avg_Tput_DL, merged_values_Avg_Tput_DL = merge_series(parsed_categories_Avg_Tput, existing_Avg_Tput[0], dates_NR_2, USER_THR_DL_data)
    merged_dates_Avg_Tput_UL, merged_values_Avg_Tput_UL = merge_series(parsed_categories_Avg_Tput, existing_Avg_Tput[1], dates_NR_2, USER_THR_UL_data)
    
    combined_values_Avg_Tput = {
    "Avg MAC user thp DL exc DRX sleep": merged_values_Avg_Tput_DL,
    "Avg MAC user thp UL": merged_values_Avg_Tput_UL} 
        
        
    update_ppt_chart(
        slide_number_5G_Traffic,
        chart_title_Avg_Tput,
        merged_dates_Avg_Tput_DL,
        combined_values_Avg_Tput,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True) 
########################## PS ##############################
#############################################################

##### Get 5G UPPER Data From Excel #################
    excel_path_2=r"C:\Users\mreem\Documents\tools\daily\PS Query 8.xlsx"
    sheet_name = 'Merge1'
    df = pd.read_excel(excel_path_2, sheet_name=sheet_name)

    # Extract data for LTE CSSR and LTE DCR
    df['Period start time'] = pd.to_datetime(df['Period start time'], errors='coerce')
    dates_NR = df['Period start time'].dt.strftime('%m-%d-%Y').tolist()  # Ensure dates are strings

    
 #### 2G KPI ####################   



    cssr_2G = df['2G Call Setup Success Rate-Speech'].tolist()
    dcr_2G = df['2G Drop Call Rate-Speech_new'].tolist()
    cssr_2G = [float(value) for value in cssr_2G]
    dcr_2G = [float(value) for value in dcr_2G]
    Traffic_2G=df['2G Busy hour Traffic - Speech'].tolist()
    Traffic_2G = [float(value) for value in Traffic_2G] 
    Traffic_Total_2G=df['2G Data'].tolist()
    Traffic_Total_2G = [float(value) for value in Traffic_Total_2G] 
    
 #### 3G KPI ####################   


    cssr_3G_V = df['3G CAll SETUP SUCCESS RATE - SPEECH'].tolist()
    cssr_3G_D = df['3G CALL SETUP SUCCESS RATE - HSDPA'].tolist()
    dcr_3G_V = df['3G VOICE_DCR'].tolist()
    dcr_3G_D = df['DCR_HSDPA'].tolist()

    cssr_3G_V = [float(value) for value in cssr_3G_V]
    dcr_3G_V = [float(value) for value in dcr_3G_V]
    cssr_3G_D = [float(value) for value in cssr_3G_D]
    dcr_3G_D = [float(value) for value in dcr_3G_D]
    Traffic_3G=df['TRAFFIC_SPEECH2'].tolist()
    Traffic_3G = [float(value) for value in Traffic_3G]
    Total_3G_Traffic = df['3G Data'].tolist()
    Total_3G_Traffic = [float(value) for value in Total_3G_Traffic]

    
    

 #### 4G KPI ####################   

    cssr_4G_D = df['4G/LTE CALL SETUP SUCCESS RATE'].tolist()
    cssr_4G_V = df['CSSR_VOLTE'].tolist()
    dcr_4G_D = df['4G/LTE DROP CALL RATE (ALL)'].tolist()
    dcr_4G_V = df['DCR_VOLTE'].tolist()  
    
    Traffic_4G_Data = df['4G Data'].tolist()   
    Traffic_4G_VO = df['VoLTE total traffic'].tolist() 
    Traffic_ORG= df['Orange Traffic'].tolist() 
    Traffic_WE= df['WE Traffic'].tolist() 
    Traffic_3G_PEN= df['3G Penetration'].tolist() 
    Traffic_VOLTE_PEN= df['VoLTE Penetration'].tolist() 
    Traffic_4G_PEN= df['4G Penetration'].tolist() 
    Traffic_WE_PEN= df['QCI9 Share'].tolist() 
    
    cssr_4G_D = [float(value) for value in cssr_4G_D]  
    cssr_4G_V = [float(value) for value in cssr_4G_V]  
    dcr_4G_D = [float(value) for value in dcr_4G_D]  
    dcr_4G_V = [float(value) for value in dcr_4G_V]     
    Traffic_4G_Data = [float(value) for value in Traffic_4G_Data]
    Traffic_4G_VO = [float(value) for value in Traffic_4G_VO]        
    Traffic_ORG = [float(value) for value in Traffic_ORG]    
    Traffic_WE = [float(value) for value in Traffic_WE]  
    Traffic_3G_PEN = [float(value) for value in Traffic_3G_PEN]      
    Traffic_VOLTE_PEN = [float(value) for value in Traffic_VOLTE_PEN]       
    Traffic_4G_PEN = [float(value) for value in Traffic_4G_PEN]       
    Traffic_WE_PEN = [float(value) for value in Traffic_WE_PEN]     

    

    Traffic_5G_DL_CO = df['5G leg'].tolist() 
    Traffic_5G_DL_CO = [float(value) for value in Traffic_5G_DL_CO]
    Traffic_x2 = df['X2 total'].tolist() 
    Traffic_x2 = [float(value) for value in Traffic_x2]
    Traffic_pdcp = df['PDCP 4G total'].tolist() 
    Traffic_pdcp = [float(value) for value in Traffic_pdcp]  
    Traffic_5GCO = df['5G Collocated'].tolist() 
    Traffic_5GCO = [float(value) for value in Traffic_5GCO]  
    
    ############################ 5G KPIs ###################
    ##### Get 5G UPPER Data From Excel #################
    shhet_5g=r"C:\Users\mreem\Documents\tools\daily\WeeklySlides_Govs-nsn2gnpo-2025_10_24-18_03_56__215 1.xlsx"
    
    sheet_name = '5G001_Govs'
    df = pd.read_excel(shhet_5g, sheet_name=sheet_name)
    filtered_df = df[df['WS_NAME'] == '5G_Colloc_5G_pilot_PortSAID']
    filtered_df['Period start time'] = pd.to_datetime(filtered_df['Period start time'], errors='coerce')
    # Extract data for LTE CSSR and LTE DCR
    dates_NR_2 = filtered_df['Period start time'].dt.strftime('%Y-%m-%d').tolist()  # Ensure dates are strings
    NR_cssr_data = filtered_df['NSA call access'].tolist()
    
    NR_HOSR_INTER_data = filtered_df['IntergNB HO SR NSA'].tolist()
    NR_HOSR_INTRA_data = filtered_df['Inafreq inaDU PSC chg tot SR'].tolist()
    MAX_USER_THR_DL_data = filtered_df['Max MAC SDU Cell Thr DL DTCH'].tolist()
    MAX_USER_THR_UL_data = filtered_df['Max MAC SDU Cell Thr UL DTCH'].tolist()
    USER_THR_UL_data = filtered_df['Avg MAC user thp UL'].tolist()
    NR_USER_AVG = filtered_df['NSA Avg nr user'].tolist()
    USER_THR_DL_data = filtered_df['Avg MAC user thp DL exc DRX sleep'].tolist()
    USER_THR_DL_data = [float(value) for value in USER_THR_DL_data]

    dates_NR_2= [datetime.strptime(date, '%Y-%m-%d').strftime('%m-%d-%Y') for date in dates_NR_2]
    # Convert float values in the lists if they are not already floats
    NR_cssr_data = [float(value) for value in NR_cssr_data]

    NR_HOSR_INTER_data = [float(value) for value in NR_HOSR_INTER_data]
    NR_HOSR_INTRA_data = [float(value) for value in NR_HOSR_INTRA_data]
    MAX_USER_THR_DL_data = [float(value) for value in MAX_USER_THR_DL_data]
    MAX_USER_THR_UL_data = [float(value) for value in MAX_USER_THR_UL_data]
    USER_THR_UL_data = [float(value) for value in USER_THR_UL_data]
    NR_USER_AVG = [float(value) for value in NR_USER_AVG]

#### 5G Rest of KPI ####################
    sheet_name_5g = '5G_DCR'
    df_5g = pd.read_excel(shhet_5g, sheet_name=sheet_name_5g)
    filtered_df_5g = df_5g[df_5g['WS_NAME'] == '5G_Colloc_5G_pilot_PortSAID']
    NR_dcr_data = filtered_df_5g['NSA SgNB t abn rel R excl X2 rst'].tolist()
    NR_dcr_data = [float(value) for value in NR_dcr_data]

########## 1st slide ########################################################################################################################################
    slide_number_Voice_Traffic = 6
    chart_title_Voice_Traffic = "Voice Traffic Evolution"
    chart_title_Data_Traffic = "Data Traffic Evolution"
    chart_title_WE_Traffic = "National Roaming Traffic Share"
    chart_title_NR_Traffic = "5G Collocated Traffic Share"

########## Update Voice_Traffic KPIS (same style as CSSR/DCR)###################
    categories_Voice_Traffic, existing_Voice_Traffic = extract_ppt_data(slide_number_Voice_Traffic, chart_title_Voice_Traffic, pptx_file)
    parsed_categories_Voice_Traffic = parse_categories(categories_Voice_Traffic)
    existing_Voice_Traffic[0] = [float(v) for v in existing_Voice_Traffic[0]] 
    existing_Voice_Traffic[1] = [float(v) for v in existing_Voice_Traffic[1]] 
    existing_Voice_Traffic[2] = [float(v) for v in existing_Voice_Traffic[2]] 
    existing_Voice_Traffic[3] = [float(v) for v in existing_Voice_Traffic[3]]
    existing_Voice_Traffic[4] = [float(v) for v in existing_Voice_Traffic[4]]
    combined_values_Voice_Traffic = {
    "2G Busy hour Traffic - Speech": existing_Voice_Traffic[0] + Traffic_2G,
    "Traffic Voice": existing_Voice_Traffic[1] + Traffic_3G,
   "VoLTE total traffic": existing_Voice_Traffic[2] + Traffic_4G_VO,
   "VoLTE Penetration": existing_Voice_Traffic[3] + Traffic_VOLTE_PEN,
   "3G Penetration": existing_Voice_Traffic[3] + Traffic_3G_PEN
   } 
        
    merged_dates_Voice_Traffic_2G, merged_values_Voice_Traffic_2G = merge_series(parsed_categories_Voice_Traffic, existing_Voice_Traffic[0], dates_NR, Traffic_2G)
    merged_dates_Voice_Traffic_3G, merged_values_Voice_Traffic_3G = merge_series(parsed_categories_Voice_Traffic, existing_Voice_Traffic[1], dates_NR, Traffic_3G)
    merged_dates_Voice_Traffic_4G, merged_values_Voice_Traffic_4G = merge_series(parsed_categories_Voice_Traffic, existing_Voice_Traffic[2], dates_NR, Traffic_4G_VO)
    merged_dates_Voice_Traffic_VP, merged_values_Voice_Traffic_VP = merge_series(parsed_categories_Voice_Traffic, existing_Voice_Traffic[3], dates_NR, Traffic_VOLTE_PEN)
    merged_dates_Voice_Traffic_3GP, merged_values_Voice_Traffic_3GP = merge_series(parsed_categories_Voice_Traffic, existing_Voice_Traffic[4], dates_NR, Traffic_3G_PEN)    
    combined_values_Voice_Traffic = {
    "2G Busy hour Traffic - Speech": merged_values_Voice_Traffic_2G,
    "Traffic Voice": merged_values_Voice_Traffic_3G,
   "VoLTE total traffic": merged_values_Voice_Traffic_4G,
   "VoLTE Penetration": merged_values_Voice_Traffic_VP,
   "3G Penetration": merged_values_Voice_Traffic_3GP} 
        
        
    update_ppt_chart(
        slide_number_Voice_Traffic,
        chart_title_Voice_Traffic,
        merged_dates_Voice_Traffic_2G,
        combined_values_Voice_Traffic,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True) 

############ Update Data_Traffic KPIS (same style as CSSR/DCR)##################
    categories_Data_Traffic, existing_Data_Traffic = extract_ppt_data(slide_number_Voice_Traffic, chart_title_Data_Traffic, pptx_file)
    parsed_categories_Data_Traffic = parse_categories(categories_Data_Traffic)
    existing_Data_Traffic[0] = [float(v) for v in existing_Data_Traffic[0]] 
    existing_Data_Traffic[1] = [float(v) for v in existing_Data_Traffic[1]] 
    existing_Data_Traffic[2] = [float(v) for v in existing_Data_Traffic[2]] 
    existing_Data_Traffic[3] = [float(v) for v in existing_Data_Traffic[3]] 
    combined_values_Data_Traffic = {
    "Sum of 2G Data": existing_Data_Traffic[0] + Traffic_Total_2G,
    "Sum of 3G Data": existing_Data_Traffic[1] + Total_3G_Traffic,
   "Sum of 4G Data": existing_Data_Traffic[2] + Traffic_4G_Data,
   "LTE Penetration": existing_Data_Traffic[3] + Traffic_4G_PEN} 
        
    merged_dates_Data_Traffic_2G_d, merged_values_Data_Traffic_2G_d = merge_series(parsed_categories_Data_Traffic, existing_Data_Traffic[0], dates_NR, Traffic_Total_2G)
    merged_dates_Data_Traffic_3G_d, merged_values_Data_Traffic_3G_d = merge_series(parsed_categories_Data_Traffic, existing_Data_Traffic[1], dates_NR, Total_3G_Traffic)
    merged_dates_Data_Traffic_4G_d, merged_values_Data_Traffic_4G_d = merge_series(parsed_categories_Data_Traffic, existing_Data_Traffic[2], dates_NR, Traffic_4G_Data)
    merged_dates_Data_Traffic_tP_d, merged_values_Data_Traffic_tP_d = merge_series(parsed_categories_Data_Traffic, existing_Data_Traffic[3], dates_NR, Traffic_4G_PEN)
  
    combined_values_Data_Traffic = {
    "Sum of 2G Data": merged_values_Data_Traffic_2G_d,
    "Sum of 3G Data": merged_values_Data_Traffic_3G_d,
   "Sum of 4G Data": merged_values_Data_Traffic_4G_d,
   "LTE Penetration": merged_values_Data_Traffic_tP_d} 

        
        
    update_ppt_chart(
        slide_number_Voice_Traffic,
        chart_title_Data_Traffic,
        merged_dates_Data_Traffic_2G_d,
        combined_values_Data_Traffic,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True) 
    
    
############ Update WE_Traffic KPIS (same style as CSSR/DCR)###################
    categories_WE_Traffic, existing_WE_Traffic = extract_ppt_data(slide_number_Voice_Traffic, chart_title_WE_Traffic, pptx_file)
    parsed_categories_WE_Traffic = parse_categories(categories_WE_Traffic)
    existing_WE_Traffic[0] = [float(v) for v in existing_WE_Traffic[0]] 
    existing_WE_Traffic[1] = [float(v) for v in existing_WE_Traffic[1]] 
    existing_WE_Traffic[2] = [float(v) for v in existing_WE_Traffic[2]] 
    combined_values_WE_Traffic = {
    "Orange Traffic (QCI8)": existing_WE_Traffic[0] + Traffic_ORG,
    "WE Traffic (QCI9)": existing_WE_Traffic[1] + Traffic_WE,
   "WE traffic Share": existing_WE_Traffic[2] + Traffic_WE_PEN} 
        
    merged_dates_WE_Traffic_ORG, merged_values_WE_Traffic_ORG = merge_series(parsed_categories_WE_Traffic, existing_WE_Traffic[0], dates_NR, Traffic_ORG)
    merged_dates_WE_Traffic_WE, merged_values_WE_Traffic_WE = merge_series(parsed_categories_WE_Traffic, existing_WE_Traffic[1], dates_NR, Traffic_WE)
    merged_dates_WE_Traffic_WP, merged_values_WE_Traffic_WP = merge_series(parsed_categories_WE_Traffic, existing_WE_Traffic[2], dates_NR, Traffic_WE_PEN)

  
    combined_values_WE_Traffic = {
    "Orange Traffic (QCI8)": merged_values_WE_Traffic_ORG,
    "WE Traffic (QCI9)": merged_values_WE_Traffic_WE,
   "WE traffic Share": merged_values_WE_Traffic_WP} 

        
        
    update_ppt_chart(
        slide_number_Voice_Traffic,
        chart_title_WE_Traffic,
        merged_dates_WE_Traffic_ORG,
        combined_values_WE_Traffic,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True) 
  
############# Update NR_Traffic KPIS (same style as CSSR/DCR)###################
    categories_NR_Traffic, existing_NR_Traffic = extract_ppt_data(slide_number_Voice_Traffic, chart_title_NR_Traffic, pptx_file)
    parsed_categories_NR_Traffic = parse_categories(categories_NR_Traffic)
    existing_NR_Traffic[0] = [float(v) for v in existing_NR_Traffic[0]] 
    existing_NR_Traffic[1] = [float(v) for v in existing_NR_Traffic[1]] 
    existing_NR_Traffic[2] = [float(v) for v in existing_NR_Traffic[2]]
    existing_NR_Traffic[3] = [float(v) for v in existing_NR_Traffic[3]]

    combined_values_NR_Traffic = {
    "PDCP 4G Total": existing_NR_Traffic[0] + Traffic_pdcp,
    "5G leg": existing_NR_Traffic[1] + Traffic_5G_DL_CO,
   "X2 total": existing_NR_Traffic[2] + Traffic_x2,
   "5G Collocated": existing_NR_Traffic[3] + Traffic_5GCO } 
        
    merged_dates_NR_Traffic_LTE, merged_values_NR_Traffic_LTE = merge_series(parsed_categories_NR_Traffic, existing_NR_Traffic[0], dates_NR, Traffic_pdcp)
    merged_dates_NR_Traffic_NR, merged_values_NR_Traffic_NR = merge_series(parsed_categories_NR_Traffic, existing_NR_Traffic[1], dates_NR, Traffic_5G_DL_CO)
    merged_dates_NR_Traffic_NR, merged_values_NR_Traffic_x2 = merge_series(parsed_categories_NR_Traffic, existing_NR_Traffic[2], dates_NR, Traffic_x2)
    merged_dates_NR_Traffic_NP, merged_values_NR_Traffic_NP = merge_series(parsed_categories_NR_Traffic, existing_NR_Traffic[3], dates_NR, Traffic_5GCO)

  
    combined_values_NR_Traffic = {
    "PDCP 4G Total": merged_values_NR_Traffic_LTE,
    "5G leg": merged_values_NR_Traffic_NR,
   "X2 total": merged_values_NR_Traffic_x2,
   "5G Collocated":merged_values_NR_Traffic_NP} 

        
        
    update_ppt_chart(
        slide_number_Voice_Traffic,
        chart_title_NR_Traffic,
        merged_dates_NR_Traffic_LTE,
        combined_values_NR_Traffic,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)   
  
########## 4th slide ##########################################################################################################################
    slide_number_LTE_Traffic = 7
    chart_title_LTE_CSSR = "4G CSSR"
    chart_title_LTE_DCR = "4G DCR"
    chart_title_LTE_Traffic = "LTE & VoLTE Traffic"

############# Update LTE_CSSR KPIS (same style as CSSR/DCR)###################
    categories_LTE_CSSR, existing_LTE_CSSR = extract_ppt_data(slide_number_LTE_Traffic, chart_title_LTE_CSSR, pptx_file)
    parsed_categories_LTE_CSSR = parse_categories(categories_LTE_CSSR)
    existing_LTE_CSSR[0] = [float(v) for v in existing_LTE_CSSR[0]] 
    existing_LTE_CSSR[1] = [float(v) for v in existing_LTE_CSSR[1]]  
    
    combined_values_LTE_CSSR = {
    "LTE_CSSR_Orange V5.2": existing_LTE_CSSR[0] + cssr_4G_D,
    "Orange CSSR VOLTE 2": existing_LTE_CSSR[1] + cssr_4G_V} 
        
    merged_dates_LTE_CSSR_D, merged_values_LTE_CSSR_D = merge_series(parsed_categories_LTE_CSSR, existing_LTE_CSSR[0], dates_NR, cssr_4G_D)
    merged_dates_LTE_CSSRc_3G, merged_values_LTE_CSSR_V = merge_series(parsed_categories_LTE_CSSR, existing_LTE_CSSR[1], dates_NR, cssr_4G_V)
    
    combined_values_LTE_CSSR = {
    "LTE_CSSR_Orange V5.2": merged_values_LTE_CSSR_D,
    "Orange CSSR VOLTE 2": merged_values_LTE_CSSR_V} 
        
        
    update_ppt_chart(
        slide_number_LTE_Traffic,
        chart_title_LTE_CSSR,
        merged_dates_LTE_CSSR_D,
        combined_values_LTE_CSSR,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)    
    
############# Update LTE_DCR KPIS (same style as CSSR/DCR)###################
    categories_LTE_DCR, existing_LTE_DCR = extract_ppt_data(slide_number_LTE_Traffic, chart_title_LTE_DCR, pptx_file)
    parsed_categories_LTE_DCR = parse_categories(categories_LTE_DCR)
    existing_LTE_DCR[0] = [float(v) for v in existing_LTE_DCR[0]] 
    existing_LTE_DCR[1] = [float(v) for v in existing_LTE_DCR[1]]  
    
    combined_values_LTE_DCR = {
    "LTE_DCR_v5_RRC_TA_filter": existing_LTE_DCR[0] + dcr_4G_D,
    "Orange VOLTE DCR V5.2": existing_LTE_DCR[1] + dcr_4G_V} 
        
    merged_dates_LTE_DCR_D, merged_values_LTE_DCR_D = merge_series(parsed_categories_LTE_DCR, existing_LTE_DCR[0], dates_NR, dcr_4G_D)
    merged_dates_LTE_DCR_V, merged_values_LTE_DCR_V = merge_series(parsed_categories_LTE_DCR, existing_LTE_DCR[1], dates_NR, dcr_4G_V)
    
    combined_values_LTE_DCR = {
    "LTE_DCR_v5_RRC_TA_filter": merged_values_LTE_DCR_D,
    "Orange VOLTE DCR V5.2": merged_values_LTE_DCR_V} 
        
        
    update_ppt_chart(
        slide_number_LTE_Traffic,
        chart_title_LTE_DCR,
        merged_dates_LTE_DCR_D,
        combined_values_LTE_DCR,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)     
    
############# Update LTE_traff KPIS (same style as CSSR/DCR)###################
    categories_LTE_traff, existing_LTE_traff = extract_ppt_data(slide_number_LTE_Traffic, chart_title_LTE_Traffic, pptx_file)
    parsed_categories_LTE_traff = parse_categories(categories_LTE_traff)
    existing_LTE_traff[0] = [float(v) for v in existing_LTE_traff[0]] 
    existing_LTE_traff[1] = [float(v) for v in existing_LTE_traff[1]]  
    
    combined_values_LTE_traff = {
    "4G Data Traffic": existing_LTE_traff[0] + Traffic_4G_Data,
    "VoLTE total traffic": existing_LTE_traff[1] + Traffic_4G_VO} 
        
    merged_dates_LTE_traff_D, merged_values_LTE_traff_D = merge_series(parsed_categories_LTE_traff, existing_LTE_traff[0], dates_NR, Traffic_4G_Data)
    merged_dates_LTE_traff_V, merged_values_LTE_traff_V = merge_series(parsed_categories_LTE_traff, existing_LTE_traff[1], dates_NR, Traffic_4G_VO)
    
    combined_values_LTE_traff = {
    "4G Data Traffic": merged_values_LTE_traff_D,
    "VoLTE total traffic": merged_values_LTE_traff_V} 
        
        
    update_ppt_chart(
        slide_number_LTE_Traffic,
        chart_title_LTE_Traffic,
        merged_dates_LTE_traff_D,
        combined_values_LTE_traff,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)    
    
########## 6th slide ##########################################################################################################################
    slide_number_3G_Traffic = 8
    chart_title_3G_CSSR = "3G CSSR"
    chart_title_3G_DCR = "3G DCR"
    chart_title_2G_CSSR = "2G CSSR"
    chart_title_2G_DCR = "2G DCR"


############# Update 3G_CSSR KPIS (same style as CSSR/DCR)###################
    categories_3G_CSSR, existing_3G_CSSR = extract_ppt_data(slide_number_3G_Traffic, chart_title_3G_CSSR, pptx_file)
    parsed_categories_3G_CSSR = parse_categories(categories_3G_CSSR)
    existing_3G_CSSR[0] = [float(v) for v in existing_3G_CSSR[0]] 
    existing_3G_CSSR[1] = [float(v) for v in existing_3G_CSSR[1]]  
    
    combined_values_3G_CSSR = {
    "3G CAll SETUP SUCCESS RATE - SPEECH": existing_3G_CSSR[0] + cssr_3G_V,
    "3G CALL SETUP SUCCESS RATE - HSDPA": existing_3G_CSSR[1] + cssr_3G_D} 
        
    merged_dates_3G_CSSR_V, merged_values_3G_CSSR_V = merge_series(parsed_categories_3G_CSSR, existing_3G_CSSR[0], dates_NR, cssr_3G_V)
    merged_dates_3G_CSSR_D, merged_values_3G_CSSR_D = merge_series(parsed_categories_3G_CSSR, existing_3G_CSSR[1], dates_NR, cssr_3G_D)
    
    combined_values_3G_CSSR = {
    "3G CAll SETUP SUCCESS RATE - SPEECH": merged_values_3G_CSSR_V,
    "3G CALL SETUP SUCCESS RATE - HSDPA": merged_values_3G_CSSR_D} 
        
        
    update_ppt_chart(
        slide_number_3G_Traffic,
        chart_title_3G_CSSR,
        merged_dates_3G_CSSR_D,
        combined_values_3G_CSSR,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)    
    

############# Update 3G_DCR KPIS (same style as CSSR/DCR)###################
    categories_3G_DCR, existing_3G_DCR = extract_ppt_data(slide_number_3G_Traffic, chart_title_3G_DCR, pptx_file)
    parsed_categories_3G_DCR = parse_categories(categories_3G_DCR)
    existing_3G_DCR[0] = [float(v) for v in existing_3G_DCR[0]] 
    existing_3G_DCR[1] = [float(v) for v in existing_3G_DCR[1]]  
    
    combined_values_3G_DCR = {
    "voicedcr3g": existing_3G_DCR[0] + dcr_3G_V,
    "HSDPA_DCR_Weekly": existing_3G_DCR[1] + dcr_3G_D} 
        
    merged_dates_3G_DCR_V, merged_values_3G_DCR_V = merge_series(parsed_categories_3G_DCR, existing_3G_DCR[0], dates_NR, dcr_3G_V)
    merged_dates_3G_DCR_D, merged_values_3G_DCR_D = merge_series(parsed_categories_3G_DCR, existing_3G_DCR[1], dates_NR, dcr_3G_D)
    
    combined_values_3G_DCR = {
    "voicedcr3g": merged_values_3G_DCR_V,
    "HSDPA_DCR_Weekly": merged_values_3G_DCR_D} 
        
        
    update_ppt_chart(
        slide_number_3G_Traffic,
        chart_title_3G_DCR,
        merged_dates_3G_DCR_D,
        combined_values_3G_DCR,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)     
    
############# Update 2G_CSSR KPIS (same style as CSSR/DCR)###################
    categories_2G_CSSR, existing_2G_CSSR = extract_ppt_data(slide_number_3G_Traffic, chart_title_2G_CSSR, pptx_file)
    parsed_categories_2G_CSSR = parse_categories(categories_2G_CSSR)
    existing_2G_CSSR[0] = [float(v) for v in existing_2G_CSSR[0]] 

    
    combined_values_2G_CSSR = {
    "GSMCSSRVoicenew": existing_2G_CSSR[0] + cssr_2G}
        
    merged_dates_2G_CSSR_V, merged_values_2G_CSSR_V = merge_series(parsed_categories_2G_CSSR, existing_2G_CSSR[0], dates_NR, cssr_2G)

    
    combined_values_2G_CSSR = {
    "GSMCSSRVoicenew": merged_values_2G_CSSR_V}
        
        
    update_ppt_chart(
        slide_number_3G_Traffic,
        chart_title_2G_CSSR,
        merged_dates_2G_CSSR_V,
        combined_values_2G_CSSR,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)    
    

############# Update 2G_DCR KPIS (same style as CSSR/DCR)###################
    categories_2G_DCR, existing_2G_DCR = extract_ppt_data(slide_number_3G_Traffic, chart_title_2G_DCR, pptx_file)
    parsed_categories_2G_DCR = parse_categories(categories_2G_DCR)
    existing_2G_DCR[0] = [float(v) for v in existing_2G_DCR[0]] 

    
    combined_values_2G_DCR = {
    "GSMCSSRVoicenew": existing_2G_DCR[0] + dcr_2G}
        
    merged_dates_2G_DCR_V, merged_values_2G_DCR_V = merge_series(parsed_categories_2G_DCR, existing_2G_DCR[0], dates_NR, dcr_2G)

    
    combined_values_2G_DCR = {
    "GSMCSSRVoicenew": merged_values_2G_DCR_V}
        
        
    update_ppt_chart(
        slide_number_3G_Traffic,
        chart_title_2G_DCR,
        merged_dates_2G_DCR_V,
        combined_values_2G_DCR,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)
########## 6th slide ##########################################################################################################################
    slide_number_5G_Traffic = 9
    chart_title_5G_CSSR = "CSSR"
    chart_title_5G_DCR = "DCR"
    chart_title_5G_HOSR = "HOSR"
    chart_title_NSA_USER = "NSA Users"
    chart_title_Avg_Tput = "Avg Tput"
    chart_title_Max_Tput = "Max Tput"
    
############# Update 5G_CSSR KPIS (same style as CSSR/DCR)###################
    categories_5G_CSSR, existing_5G_CSSR = extract_ppt_data(slide_number_5G_Traffic, chart_title_5G_CSSR, pptx_file)
    parsed_categories_5G_CSSR = parse_categories(categories_5G_CSSR)
    existing_5G_CSSR[0] = [float(v) for v in existing_5G_CSSR[0]] 
    print(categories_5G_CSSR)
    
    combined_values_5G_CSSR = {
    "NSA call access": existing_5G_CSSR[0] + NR_cssr_data}
        
    merged_dates_5G_CSSR_V, merged_values_5G_CSSR_V = merge_series(parsed_categories_5G_CSSR, existing_5G_CSSR[0], dates_NR_2, NR_cssr_data)

    
    combined_values_5G_CSSR = {
    "NSA call access": merged_values_5G_CSSR_V}
        
        
    update_ppt_chart(
        slide_number_5G_Traffic,
        chart_title_5G_CSSR,
        merged_dates_5G_CSSR_V,
        combined_values_5G_CSSR,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)    
  

############# Update 5G_DCR KPIS (same style as CSSR/DCR)###################
    categories_5G_DCR, existing_5G_DCR = extract_ppt_data(slide_number_5G_Traffic, chart_title_5G_DCR, pptx_file)
    parsed_categories_5G_DCR = parse_categories(categories_5G_DCR)
    existing_5G_DCR[0] = [float(v) for v in existing_5G_DCR[0]] 

    
    combined_values_5G_DCR = {
    "5G NSA DROP CALL RATE NRBTS": existing_5G_DCR[0] + NR_dcr_data}
        
    merged_dates_5G_DCR_V, merged_values_5G_DCR_V = merge_series(parsed_categories_5G_DCR, existing_5G_DCR[0], dates_NR_2, NR_dcr_data)

    
    combined_values_5G_DCR = {
    "5G NSA DROP CALL RATE NRBTS": merged_values_5G_DCR_V}
        
        
    update_ppt_chart(
        slide_number_5G_Traffic,
        chart_title_5G_DCR,
        merged_dates_5G_DCR_V,
        combined_values_5G_DCR,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)    
    
############# Update 5G_DCR KPIS (same style as CSSR/DCR)###################
    categories_NSA_USER, existing_NSA_USER = extract_ppt_data(slide_number_5G_Traffic, chart_title_NSA_USER, pptx_file)
    parsed_categories_NSA_USER = parse_categories(categories_NSA_USER)
    existing_NSA_USER[0] = [float(v) for v in existing_NSA_USER[0]] 

    
    combined_values_NSA_USER = {
    "NSA Avg nr user": existing_NSA_USER[0] + NR_USER_AVG}
        
    merged_dates_NSA_USER_V, merged_values_NSA_USER_V = merge_series(parsed_categories_NSA_USER, existing_NSA_USER[0], dates_NR_2, NR_USER_AVG)

    
    combined_values_NSA_USER = {
    "NSA Avg nr user": merged_values_NSA_USER_V}
        
 
    update_ppt_chart(
        slide_number_5G_Traffic,
        chart_title_NSA_USER,
        merged_dates_NSA_USER_V,
        combined_values_NSA_USER,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)    
    
############# Update HOSR KPIS (same style as CSSR/DCR)###################
    categories_NR_HOSR, existing_NR_HOSR = extract_ppt_data(slide_number_5G_Traffic, chart_title_5G_HOSR, pptx_file)
    parsed_categories_NR_HOSR = parse_categories(categories_NR_HOSR)
    existing_NR_HOSR[0] = [float(v) for v in existing_NR_HOSR[0]] 
    existing_NR_HOSR[1] = [float(v) for v in existing_NR_HOSR[1]]  
    
    combined_values_NR_HOSR = {
    "IntergNB HO SR NSA": existing_NR_HOSR[0] + NR_HOSR_INTER_data,
    "Inafreq inaDU PSC chg tot SR": existing_NR_HOSR[1] + NR_HOSR_INTRA_data} 
        
    merged_dates_NR_HOSR_inter, merged_values_NR_HOSR_inter = merge_series(parsed_categories_NR_HOSR, existing_NR_HOSR[0], dates_NR_2, NR_HOSR_INTER_data)
    merged_dates_NR_HOSR_intra, merged_values_NR_HOSR_intra = merge_series(parsed_categories_NR_HOSR, existing_NR_HOSR[1], dates_NR_2, NR_HOSR_INTRA_data)
    
    combined_values_NR_HOSR = {
    "IntergNB HO SR NSA": merged_values_NR_HOSR_inter,
    "Inafreq inaDU PSC chg tot SR": merged_values_NR_HOSR_intra} 
        
        
    update_ppt_chart(
        slide_number_5G_Traffic,
        chart_title_5G_HOSR,
        merged_dates_NR_HOSR_inter,
        combined_values_NR_HOSR,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)     

############# Update Max Tput KPIS (same style as CSSR/DCR)###################

    categories_max_tput, existing_max_tput = extract_ppt_data(slide_number_5G_Traffic, chart_title_Max_Tput, pptx_file)
    parsed_categories_max_tput = parse_categories(categories_max_tput)
    existing_max_tput[0] = [float(v) for v in existing_max_tput[0]] 
    existing_max_tput[1] = [float(v) for v in existing_max_tput[1]]  
    
    combined_values_max_tput = {
    "Max MAC SDU Cell Thr DL DTCH": existing_max_tput[0] + MAX_USER_THR_DL_data,
    "Max MAC SDU Cell Thr UL DTCH": existing_max_tput[1] + MAX_USER_THR_UL_data} 
        
    merged_dates_max_tput_DL, merged_values_max_tput_DL = merge_series(parsed_categories_max_tput, existing_max_tput[0], dates_NR_2, MAX_USER_THR_DL_data)
    merged_dates_max_tput_UL, merged_values_max_tput_UL = merge_series(parsed_categories_max_tput, existing_max_tput[1], dates_NR_2, MAX_USER_THR_UL_data)
    
    combined_values_max_tput = {
    "Max MAC SDU Cell Thr DL DTCH": merged_values_max_tput_DL,
    "Max MAC SDU Cell Thr UL DTCH": merged_values_max_tput_UL} 
        
        
    update_ppt_chart(
        slide_number_5G_Traffic,
        chart_title_Max_Tput,
        merged_dates_max_tput_DL,
        combined_values_max_tput,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True)   
    
############# Update Avg Tput KPIS (same style as CSSR/DCR)###################

    categories_Avg_Tput, existing_Avg_Tput = extract_ppt_data(slide_number_5G_Traffic, chart_title_Avg_Tput, pptx_file)
    parsed_categories_Avg_Tput = parse_categories(categories_Avg_Tput)
    existing_Avg_Tput[0] = [float(v) for v in existing_Avg_Tput[0]] 
    existing_Avg_Tput[1] = [float(v) for v in existing_Avg_Tput[1]]  
    
    combined_values_Avg_Tput = {
    "Avg MAC user thp DL exc DRX sleep": existing_Avg_Tput[0] + USER_THR_DL_data,
    "Avg MAC user thp UL": existing_Avg_Tput[1] + USER_THR_UL_data} 
        
    merged_dates_Avg_Tput_DL, merged_values_Avg_Tput_DL = merge_series(parsed_categories_Avg_Tput, existing_Avg_Tput[0], dates_NR_2, USER_THR_DL_data)
    merged_dates_Avg_Tput_UL, merged_values_Avg_Tput_UL = merge_series(parsed_categories_Avg_Tput, existing_Avg_Tput[1], dates_NR_2, USER_THR_UL_data)
    
    combined_values_Avg_Tput = {
    "Avg MAC user thp DL exc DRX sleep": merged_values_Avg_Tput_DL,
    "Avg MAC user thp UL": merged_values_Avg_Tput_UL} 
        
        
    update_ppt_chart(
        slide_number_5G_Traffic,
        chart_title_Avg_Tput,
        merged_dates_Avg_Tput_DL,
        combined_values_Avg_Tput,
        series_name=None,  # not used for multi-series
        ppt_file=pptx_file,
        multi_series=True) 
    print("Successfully updated PowerPoint slide with aggregated data.")
if __name__ == "__main__":
    main()    
######################### Port Saaed###########################################

###############################################################################
